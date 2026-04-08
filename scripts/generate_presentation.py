#!/usr/bin/env python3
"""Generate CodeStrike presentation for COMP 4210 Ethical Hacking."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors matching dashboard theme
BG_DARK = RGBColor(0x09, 0x09, 0x0B)       # zinc-950
BG_CARD = RGBColor(0x18, 0x18, 0x1B)       # zinc-900
BORDER = RGBColor(0x27, 0x27, 0x2A)        # zinc-800
TEXT_WHITE = RGBColor(0xF4, 0xF4, 0xF5)    # zinc-100
TEXT_LIGHT = RGBColor(0xA1, 0xA1, 0xAA)    # zinc-400
TEXT_DIM = RGBColor(0x71, 0x71, 0x7A)      # zinc-500
EMERALD = RGBColor(0x34, 0xD3, 0x99)       # emerald-400
RED = RGBColor(0xF8, 0x71, 0x71)           # red-400
AMBER = RGBColor(0xFB, 0xBF, 0x24)         # amber-400
BLUE = RGBColor(0x60, 0xA5, 0xFA)          # blue-400
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)        # purple-400
TEAL = RGBColor(0x2D, 0xD4, 0xBF)          # teal-400

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_LIGHT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)

        if isinstance(item, tuple):
            # (text, color) tuple
            p.text = item[0]
            p.font.color.rgb = item[1]
        else:
            p.text = item
            p.font.color.rgb = color
        p.font.size = Pt(font_size)
        p.font.name = "Calibri"
    return txBox

def add_card(slide, left, top, width, height, title, value, subtitle="", value_color=EMERALD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    add_text(slide, left + 0.2, top + 0.15, width - 0.4, 0.3, title, font_size=10, color=TEXT_DIM, bold=False)
    add_text(slide, left + 0.2, top + 0.45, width - 0.4, 0.5, value, font_size=28, color=value_color, bold=True)
    if subtitle:
        add_text(slide, left + 0.2, top + 1.0, width - 0.4, 0.3, subtitle, font_size=10, color=TEXT_DIM)

def add_speaker_tag(slide, speaker_name, color=TEAL):
    add_text(slide, 11.5, 0.2, 1.5, 0.3, f"Speaker: {speaker_name}", font_size=9, color=color)

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1, 1.5, 11, 1.2, "{/} CodeStrike", font_size=52, color=EMERALD, bold=True)
add_text(slide, 1, 2.8, 11, 0.8, "Security Benchmark for AI-Generated Code", font_size=28, color=TEXT_WHITE)
add_text(slide, 1, 3.8, 11, 0.5, "Can AI Write Secure Web Applications?", font_size=20, color=AMBER)
add_text(slide, 1, 5.0, 11, 0.4, "COMP 4210 Ethical Hacking  |  Group 8  |  April 2026", font_size=16, color=TEXT_DIM)
add_text(slide, 1, 5.5, 11, 0.4, "Yassh  |  Alex  |  Jordan  |  [Member 4]", font_size=14, color=TEXT_DIM)
add_notes(slide, """YASSH (Speaker 1):
"Good morning everyone. Today we're presenting CodeStrike — a security benchmark that answers one critical question: When AI writes code, is that code secure?

We all know AI code generation is exploding — GitHub Copilot, ChatGPT, Claude. Developers are using AI to write production code every day. But nobody has systematically tested whether that code is actually safe.

That's what CodeStrike does. We took 15 AI models, asked each one to build web applications, then attacked those applications with real exploit payloads to find vulnerabilities. What we found is... not great. Let me walk you through it."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1, 0.5, 11, 0.6, "Agenda", font_size=36, color=TEXT_WHITE, bold=True)

agenda_items = [
    "1. Scope & Objectives — What we're testing and why",
    "2. Reconnaissance — Understanding AI-generated code as attack surface",
    "3. Scanning & Enumeration — 3-layer detection approach",
    "4. Exploitation — Attack vectors and results",
    "5. Post-Exploitation — What we found inside vulnerable apps",
    "6. Summary & Recommendations — Key findings and mitigations",
    "7. Live Dashboard Demo — Interactive walkthrough of results",
]
add_bullet_list(slide, 1, 1.5, 10, 5, agenda_items, font_size=20, color=TEXT_LIGHT)

speakers = "Yassh: Scope & Objectives  |  Alex: Recon & Scanning  |  Jordan: Exploitation  |  [Member 4]: Summary & Demo"
add_text(slide, 1, 6.5, 11, 0.4, speakers, font_size=12, color=TEXT_DIM)

add_notes(slide, """YASSH:
"Here's our agenda. We'll cover the full penetration testing lifecycle applied to AI-generated code. I'll start with scope and objectives, Alex will cover reconnaissance and scanning, Jordan handles exploitation, and [Member 4] wraps up with summary and our live dashboard demo.

The unique thing about our project is that our 'target' isn't one application — it's 4,505 applications generated by 15 different AI models. Let's dive in."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 3: The Problem
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Yassh")
add_text(slide, 1, 0.5, 11, 0.6, "The Problem", font_size=36, color=TEXT_WHITE, bold=True)
add_text(slide, 1, 1.3, 11, 0.5, "AI is writing production code. Nobody is checking if it's secure.", font_size=20, color=AMBER)

add_card(slide, 1, 2.2, 3.5, 1.4, "AI CODE ADOPTION", "97%", "of developers use AI coding tools (GitHub 2025)", BLUE)
add_card(slide, 5, 2.2, 3.5, 1.4, "SECURITY REVIEWS", "< 30%", "of AI-generated code gets security review", RED)
add_card(slide, 9, 2.2, 3.5, 1.4, "OUR QUESTION", "Is it safe?", "Systematic testing at scale", EMERALD)

add_bullet_list(slide, 1, 4.2, 11, 2.5, [
    "AI models generate code that compiles and runs — but is it secure?",
    "Existing benchmarks test functionality, not security",
    "BaxBench (ICML 2025) was the first security benchmark — but limited",
    "We extended it significantly and named it CodeStrike",
], font_size=18)

add_notes(slide, """YASSH:
"Here's the problem. 97% of developers now use AI coding tools. But less than 30% of AI-generated code gets a proper security review. Everyone assumes if the code works, it's fine. But working code and secure code are very different things.

BaxBench from ICML 2025 was the first attempt at systematically testing AI code security. We took it as our foundation, extended it significantly — more scenarios, more CWEs, more exploit vectors, more models — and created CodeStrike."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 4: Scope — Domain & Contribution
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Yassh")
add_text(slide, 1, 0.5, 11, 0.6, "Scope: Domain & Our Contribution", font_size=36, color=TEXT_WHITE, bold=True)

add_text(slide, 1, 1.3, 5, 0.4, "Domain: Web Application Security", font_size=20, color=EMERALD, bold=True)
add_bullet_list(slide, 1, 1.9, 5, 2, [
    "AI-generated REST API backends",
    "3 frameworks: Flask, Express, Go-Fiber",
    "3 safety prompt levels: none, generic, specific",
    "35 unique web application scenarios",
], font_size=15)

add_text(slide, 7, 1.3, 5.5, 0.4, "What We Extended (vs Original BaxBench)", font_size=20, color=PURPLE, bold=True)
add_bullet_list(slide, 7, 1.9, 5.5, 4, [
    ("Scenarios: 28 -> 35 (+7 OWASP 2025)", EMERALD),
    ("CWEs monitored: 23 -> 38 (+15 new types)", EMERALD),
    ("XSS vectors: 2 -> 25 (polyglot, encoding bypasses)", EMERALD),
    ("SQLi vectors: 8 -> 25 (UNION, time-based blind)", EMERALD),
    ("Added: SSRF, IDOR, XXE, mass assignment scenarios", EMERALD),
    ("Added: SAST module (weak crypto, hardcoded creds)", EMERALD),
    ("Added: 4 universal security tests per scenario", EMERALD),
    ("Fixed: 'secure by crash' measurement bug", RED),
], font_size=14)

add_notes(slide, """YASSH:
"Our domain is web application security — specifically, REST API backends generated by AI. We test in three frameworks: Python Flask, JavaScript Express, and Go Fiber.

Our unique contribution over the original BaxBench is significant. We added 7 new scenarios targeting OWASP Top 10 2025 — things like SSRF, IDOR, XXE that the original didn't test. We expanded from 23 to 38 CWE types. We went from 2 XSS test vectors to 25 — including polyglot payloads and encoding bypasses.

And critically, we discovered and fixed a measurement bug in the original benchmark where crashed code was counted as 'secure.' We'll show you how this inflated security rates from 94% down to the real number of 4.4%."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 5: Scope — OWASP Top 10 2025
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Yassh")
add_text(slide, 1, 0.5, 11, 0.6, "Vulnerabilities Tested: OWASP Top 10 2025", font_size=36, color=TEXT_WHITE, bold=True)

owasp_items = [
    ("A01 Broken Access Control — CWE-284, 352, 639, 862, 863, 918", RED),
    ("A02 Cryptographic Failures — CWE-327, 338, 522, 798", RED),
    ("A03 Injection — CWE-20, 22, 78, 79, 89, 94, 117, 611", RED),
    ("A04 Insecure Design — CWE-640, 840", RED),
    ("A05 Security Misconfiguration — CWE-614, 693, 942, 1275", RED),
    ("A06 Vulnerable Components — CWE-1104", TEXT_DIM),
    ("A07 Auth Failures — CWE-287, 307, 384, 613", RED),
    ("A08 Integrity Failures — CWE-345, 347, 502", TEXT_DIM),
    ("A09 Logging Failures — CWE-209", TEXT_DIM),
    ("A10 Exceptional Conditions — CWE-400, 636, 703 (new in 2025)", RED),
]
add_bullet_list(slide, 1, 1.4, 11, 5.5, owasp_items, font_size=16)

add_text(slide, 1, 6.5, 11, 0.4, "Result: Found vulnerabilities in 7 of 10 categories  |  38 CWEs monitored  |  18 CWE types found  |  2,295 total occurrences", font_size=13, color=AMBER)

add_notes(slide, """YASSH:
"We mapped all 38 CWEs to the OWASP Top 10 2025 — the latest revision. We cover all 10 categories. The red items are categories where we found vulnerabilities. The grey ones came back clean — but two of those couldn't be properly tested because the security tests crashed 70% of the time.

In total: 38 CWE types monitored, 18 actually found, 2,295 total vulnerability occurrences across 4,505 tests. The dominant failure is A05 Security Misconfiguration with 1,456 findings — almost all missing security headers."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 6: Scope — Goals & Scale
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Yassh")
add_text(slide, 1, 0.5, 11, 0.6, "Goals & Scale of Testing", font_size=36, color=TEXT_WHITE, bold=True)

add_card(slide, 1, 1.3, 2.5, 1.4, "TOTAL TESTS", "4,505", "apps generated & tested", EMERALD)
add_card(slide, 3.8, 1.3, 2.5, 1.4, "AI MODELS", "15", "6 thinking + 9 standard", BLUE)
add_card(slide, 6.6, 1.3, 2.5, 1.4, "SCENARIOS", "35", "web app types", PURPLE)
add_card(slide, 9.4, 1.3, 2.5, 1.4, "CWEs", "38", "vulnerability types", RED)

add_text(slide, 1, 3.2, 11, 0.4, "Pentest Goals:", font_size=20, color=EMERALD, bold=True)
add_bullet_list(slide, 1, 3.8, 5.5, 3, [
    "Test if AI-generated code is vulnerable to injection attacks (SQLi, XSS, command injection)",
    "Check for authentication/authorization bypasses (IDOR, brute force, missing auth)",
    "Verify session management (CSRF, cookie flags, session fixation)",
    "Test for data exposure (hardcoded secrets, weak crypto, error leakage)",
    "Validate with industry scanner (OWASP ZAP) and manual pentesting",
], font_size=15)

add_text(slide, 7, 3.2, 5.5, 0.4, "3-Layer Validation:", font_size=20, color=AMBER, bold=True)
add_bullet_list(slide, 7, 3.8, 5.5, 3, [
    ("Layer 1: CodeStrike automated (4,505 apps)", EMERALD),
    ("Layer 2: OWASP ZAP scanning (50 apps)", AMBER),
    ("Layer 3: Manual pentesting (10 apps)", BLUE),
    ("Result: 100% precision, 27% recall", TEXT_WHITE),
    ("ZAP agreement: 14.3%", TEXT_WHITE),
], font_size=15)

add_notes(slide, """YASSH:
"Our goals were comprehensive. We wanted to test every major vulnerability category against AI code at scale. 4,505 tests across 15 models, 35 scenarios, and 38 CWE types.

But automated testing alone isn't enough. So we validated our results with three layers: CodeStrike's automated tests on all 4,505 apps, OWASP ZAP scanning on 50 selected apps, and manual penetration testing on 10 apps by our team members Alex and Jordan.

The key validation metric: CodeStrike has 100% precision — every vulnerability it reports is real. Zero false positives. Now I'll hand over to Alex for reconnaissance."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 7: Reconnaissance — Understanding the Target
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Alex")
add_text(slide, 1, 0.5, 11, 0.6, "Reconnaissance: Understanding the Target", font_size=36, color=TEXT_WHITE, bold=True)
add_text(slide, 1, 1.3, 11, 0.5, "Our 'target' is unique: AI-generated web applications running in Docker containers", font_size=18, color=AMBER)

add_text(slide, 1, 2.2, 5.5, 0.4, "Passive Reconnaissance", font_size=20, color=BLUE, bold=True)
add_bullet_list(slide, 1, 2.8, 5.5, 3, [
    "Analyzed AI model training data patterns",
    "Studied which security patterns models learn vs miss",
    "Reviewed OpenAPI specs for each scenario",
    "Identified that models prioritize functionality over security",
    "Found models rarely add security headers unless asked",
], font_size=15)

add_text(slide, 7, 2.2, 5.5, 0.4, "Active Reconnaissance", font_size=20, color=EMERALD, bold=True)
add_bullet_list(slide, 7, 2.8, 5.5, 3, [
    "Each app runs in isolated Docker container on port 5000",
    "HTTP API endpoints exposed (REST JSON APIs)",
    "Frameworks: Flask (Python), Express (Node.js), Fiber (Go)",
    "Attack surface: user input via JSON body, URL params, headers",
    "Authentication: JWT tokens, session cookies (when implemented)",
], font_size=15)

add_notes(slide, """ALEX (Speaker 2):
"Thanks Yassh. Now let me walk you through our reconnaissance phase.

Our target is unique — instead of one application, we're testing thousands of AI-generated apps. Each runs in an isolated Docker container exposing a REST API on port 5000.

For passive recon, we studied how AI models handle security. The key finding: models prioritize making code work, not making it secure. They'll implement every API endpoint correctly but forget to add security headers, CSRF tokens, or rate limiting. This informed our attack strategy.

For active recon, each app exposes JSON API endpoints. The attack surface is user input through JSON bodies, URL parameters, and headers. Some apps implement JWT auth, some use session cookies, and some have no auth at all — that's already a finding."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 8: Reconnaissance — Attack Surface
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Alex")
add_text(slide, 1, 0.5, 11, 0.6, "Attack Surface: 35 Scenarios", font_size=36, color=TEXT_WHITE, bold=True)

add_text(slide, 1, 1.3, 5.5, 0.4, "Original 28 Scenarios", font_size=18, color=TEXT_LIGHT, bold=True)
add_bullet_list(slide, 1, 1.8, 5.5, 2.5, [
    "Login, Forum, Wiki — auth & session targets",
    "Calculator, Compiler — code injection targets",
    "FileSearch, ImageTransfer — path traversal targets",
    "Monitor, PDFCat, SongDownloader — OS injection targets",
    "SecretStorage, ShopOverview — credential targets",
], font_size=14)

add_text(slide, 7, 1.3, 5.5, 0.4, "7 New OWASP 2025 Scenarios (Our Addition)", font_size=18, color=EMERALD, bold=True)
add_bullet_list(slide, 7, 1.8, 5.5, 3, [
    ("AdminPanel — IDOR, privilege escalation (A01)", EMERALD),
    ("LinkPreview — SSRF, internal network access (A01)", EMERALD),
    ("MultiUserNotes — IDOR, broken object-level auth (A01)", EMERALD),
    ("XMLImporter — XXE injection (A05)", EMERALD),
    ("PasswordReset — weak recovery mechanism (A07)", EMERALD),
    ("UserSettings — mass assignment (A01)", EMERALD),
    ("Checkout — business logic flaws (A04)", EMERALD),
], font_size=14)

add_text(slide, 1, 5.0, 11, 0.5, "Each scenario x 3 frameworks x 3 safety prompts = 9 tests per model per scenario", font_size=16, color=AMBER)

add_notes(slide, """ALEX:
"Each scenario represents a different type of web application with specific vulnerability targets. The original BaxBench had 28 — covering injection, auth, and file handling.

We added 7 new scenarios specifically targeting OWASP Top 10 2025 gaps. AdminPanel tests for privilege escalation. LinkPreview tests SSRF — can the app be tricked into accessing internal network resources? XMLImporter tests XXE. PasswordReset tests for weak recovery flows.

Each scenario is tested 9 ways: 3 frameworks times 3 safety prompt levels. That's how we get to 4,505 total tests across 15 models."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 9: Tools Used
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Alex")
add_text(slide, 1, 0.5, 11, 0.6, "Tools & Methodology", font_size=36, color=TEXT_WHITE, bold=True)

add_text(slide, 1, 1.3, 3.5, 0.4, "CodeStrike (Custom)", font_size=18, color=EMERALD, bold=True)
add_bullet_list(slide, 1, 1.8, 3.5, 2, [
    "Python-based security testing framework",
    "25 XSS vectors (polyglot, SVG, encoding)",
    "25 SQLi vectors (UNION, blind, stacked)",
    "10 SSRF vectors, 5 blind OS injection",
    "SAST module (regex-based code analysis)",
    "Runs INSIDE Docker alongside the app",
], font_size=13)

add_text(slide, 5, 1.3, 3.5, 0.4, "OWASP ZAP", font_size=18, color=AMBER, bold=True)
add_bullet_list(slide, 5, 1.8, 3.5, 2, [
    "Industry-standard DAST scanner",
    "Tested 3 modes: baseline, full, API scan",
    "56 active scan rules available",
    "Used for external validation",
    "Ran against 50 apps in Docker",
], font_size=13)

add_text(slide, 9, 1.3, 3.5, 0.4, "Manual Pentesting", font_size=18, color=BLUE, bold=True)
add_bullet_list(slide, 9, 1.8, 3.5, 2, [
    "OWASP WSTG v4.2 methodology",
    "PTES framework",
    "Burp Suite / curl / custom scripts",
    "10 apps, 2 testers (Alex + Jordan)",
    "107-point security checklist",
], font_size=13)

add_text(slide, 1, 4.5, 11, 0.4, "Infrastructure", font_size=18, color=PURPLE, bold=True)
add_bullet_list(slide, 1, 5.0, 11, 2, [
    "Docker containers (python:3.12-bookworm, node:22-bookworm, golang:1.24-bookworm)",
    "Next.js dashboard deployed on Vercel  |  SQLite database  |  15 AI models via API (Anthropic Claude, DeepSeek, Meta Llama)",
], font_size=14)

add_notes(slide, """ALEX:
"Our toolkit has three layers. First, CodeStrike — our custom Python framework that runs inside the Docker container alongside the app. This is crucial because it can query the database directly, send authenticated requests, and test things like rate limiting that external scanners can't.

Second, OWASP ZAP — the industry standard. We tested it in three modes against 50 apps. We'll show you why it only agreed with 14% of our findings.

Third, manual pentesting. Jordan and I manually hacked 10 apps using OWASP WSTG methodology with a 107-point checklist. This is our ground truth for measuring accuracy.

All apps run in isolated Docker containers to ensure reproducibility."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 10: Scanning — 3-Layer Approach
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Alex")
add_text(slide, 1, 0.5, 11, 0.6, "Scanning & Enumeration: 3-Layer Approach", font_size=36, color=TEXT_WHITE, bold=True)

# Layer 1
add_text(slide, 1, 1.5, 3.5, 0.4, "Layer 1: CodeStrike Automated", font_size=18, color=EMERALD, bold=True)
add_text(slide, 1, 2.0, 3.5, 0.3, "4,505 apps  |  38 CWE types", font_size=14, color=TEXT_DIM)
add_bullet_list(slide, 1, 2.4, 3.5, 2.5, [
    "Sends exploit payloads to every endpoint",
    "Checks response for vulnerability indicators",
    "Accesses Docker DB for credential checks",
    "Tests rate limiting (15 rapid requests)",
    "Creates users to test IDOR/auth bypass",
], font_size=13)

# Layer 2
add_text(slide, 5, 1.5, 3.5, 0.4, "Layer 2: OWASP ZAP Scan", font_size=18, color=AMBER, bold=True)
add_text(slide, 5, 2.0, 3.5, 0.3, "50 apps  |  3 scan modes", font_size=14, color=TEXT_DIM)
add_bullet_list(slide, 5, 2.4, 3.5, 2.5, [
    "Baseline: passive header checks",
    "Full scan: all 56 active rules",
    "API scan: OpenAPI spec imported",
    "Tests from outside the container",
    "Industry standard validation",
], font_size=13)

# Layer 3
add_text(slide, 9, 1.5, 3.5, 0.4, "Layer 3: Manual Pentest", font_size=18, color=BLUE, bold=True)
add_text(slide, 9, 2.0, 3.5, 0.3, "10 apps  |  47 findings", font_size=14, color=TEXT_DIM)
add_bullet_list(slide, 9, 2.4, 3.5, 2.5, [
    "OWASP WSTG v4.2 checklist (107 checks)",
    "Tested auth, session, injection, logic",
    "Found 12 unique CWE types",
    "2 critical, 12 high, 22 medium, 11 low",
    "Ground truth for accuracy measurement",
], font_size=13)

add_text(slide, 1, 5.5, 11, 0.8, "Key Result: Manual found 47 vulns  |  CodeStrike confirmed 11 (100% precision)  |  ZAP found 2 CWE types (14.3% agreement)", font_size=16, color=AMBER)

add_notes(slide, """ALEX:
"Our scanning approach uses three layers, each adding depth.

Layer 1 is CodeStrike automated — 4,505 apps tested with 38 CWE types. It runs inside the Docker container so it can do things like query the SQLite database to check if passwords are hashed properly.

Layer 2 is OWASP ZAP — the industry standard. We ran it against 50 apps in three modes. Key finding: ZAP only confirmed 14.3% of our results, not because we're wrong, but because ZAP can't access the database or understand auth logic.

Layer 3 is manual pentesting — Alex and I tested 10 apps by hand. This found 47 vulnerabilities across 12 CWE types. This is our ground truth.

The validation: CodeStrike has 100% precision — everything it flags is real. Now Jordan will show you the exploitation phase."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 11: Exploitation — The Security Funnel
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Jordan")
add_text(slide, 1, 0.5, 11, 0.6, "Exploitation Results: The Security Funnel", font_size=36, color=TEXT_WHITE, bold=True)

# Funnel visualization
add_card(slide, 1, 1.5, 11, 1.0, "4,505 APPS GENERATED BY AI", "100%", "", BLUE)
add_card(slide, 2, 2.8, 9, 1.0, "3,217 CRASHED (71.4%)", "Code doesn't work", "syntax errors, missing imports, runtime failures", RED)
add_card(slide, 3, 4.1, 7, 1.0, "1,288 FUNCTIONAL (28.6%)", "Code works", "passes all functional tests", AMBER)
add_card(slide, 4, 5.4, 5, 1.0, "1,142 VULNERABLE (88.7% of working)", "Has CWEs", "", RED)
add_card(slide, 5, 6.0, 3, 0.7, "146 SECURE", "3.2%", "", EMERALD)

add_notes(slide, """JORDAN (Speaker 3):
"Thanks Alex. Now let me show you what happens when we actually attack these apps.

This is the security funnel — the most important slide in our presentation. We started with 4,505 AI-generated apps. Here's what happened:

71.4% — that's 3,217 apps — crashed. The AI couldn't even write working code. Syntax errors, missing imports, runtime failures. Most AI code doesn't work at all.

Of the 1,288 that survived, 88.7% — that's 1,142 — have at least one security vulnerability. We exploited them.

Only 146 apps — 3.2% of the total — are both functional AND have zero detected vulnerabilities. And of those, only 31 are truly clean when we require all security tests to have run without crashing.

Less than 1% of AI-generated code is truly secure."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 12: Exploitation — Attack Vectors
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Jordan")
add_text(slide, 1, 0.5, 11, 0.6, "Exploitation: Attack Vectors Used", font_size=36, color=TEXT_WHITE, bold=True)

add_text(slide, 1, 1.3, 5.5, 0.4, "XSS — 25 Vectors", font_size=18, color=RED, bold=True)
add_bullet_list(slide, 1, 1.8, 5.5, 1.5, [
    '<script>alert(1)</script>  (basic)',
    '<img src=x onerror=alert(1)>  (event handler)',
    '<svg/onload=alert(1)>  (SVG injection)',
    'javascript:alert(1)  (protocol handler)',
    '%3Cscript%3Ealert(1)%3C/script%3E  (URL encoded)',
], font_size=13)

add_text(slide, 7, 1.3, 5.5, 0.4, "SQL Injection — 25 Vectors", font_size=18, color=RED, bold=True)
add_bullet_list(slide, 7, 1.8, 5.5, 1.5, [
    "' OR '1'='1  (classic bypass)",
    "' UNION SELECT * FROM users--  (data extraction)",
    "'; DROP TABLE users--  (destructive)",
    "' AND SLEEP(5)--  (time-based blind)",
    "admin'/*  (comment injection)",
], font_size=13)

add_text(slide, 1, 3.8, 5.5, 0.4, "OS Command Injection", font_size=18, color=RED, bold=True)
add_bullet_list(slide, 1, 4.3, 5.5, 1.2, [
    "; ls /etc/passwd  (basic)",
    "| cat /etc/shadow  (pipe)",
    "`sleep 5`  (time-based blind)",
    "$(whoami)  (subshell)",
], font_size=13)

add_text(slide, 7, 3.8, 5.5, 0.4, "Other Attacks", font_size=18, color=RED, bold=True)
add_bullet_list(slide, 7, 4.3, 5.5, 1.5, [
    "SSRF: http://169.254.169.254/metadata (cloud)",
    "Path traversal: ../../etc/passwd (22+ variants)",
    "CSRF: missing token validation",
    "Brute force: 15 rapid login attempts",
    "IDOR: swap user A's token for user B's data",
], font_size=13)

add_notes(slide, """JORDAN:
"Here are the actual attack vectors we used. For XSS alone, we have 25 different payloads — from basic script tags to SVG injection, event handlers, and encoding bypasses. The original benchmark only had 2.

For SQL injection, another 25 vectors including UNION-based extraction, time-based blind detection, and comment injection. Again, the original had only 8.

We also test OS command injection with blind time-based detection, SSRF with cloud metadata endpoints, path traversal with 22 encoding bypass variants, and brute force with rapid request testing.

Each of these runs against every endpoint of every app. That's why we find things that ZAP misses — we're sending targeted payloads in the correct format for each API."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 13: Exploitation — CWE Results
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Jordan")
add_text(slide, 1, 0.5, 11, 0.6, "Exploitation Results: What We Found", font_size=36, color=TEXT_WHITE, bold=True)

cwe_data = [
    ("CWE-693", "Missing Headers", "1,427", "62.2%"),
    ("CWE-352", "CSRF Missing", "271", "11.8%"),
    ("CWE-307", "No Rate Limiting", "144", "6.3%"),
    ("CWE-79", "XSS", "116", "5.1%"),
    ("CWE-400", "Resource Exhaustion", "115", "5.0%"),
    ("CWE-522", "Weak Credentials", "62", "2.7%"),
    ("CWE-89", "SQL Injection", "1", "0.04%"),
]

y = 1.3
add_text(slide, 1, y, 2, 0.3, "CWE", font_size=12, color=TEXT_DIM, bold=True)
add_text(slide, 3, y, 3, 0.3, "Vulnerability", font_size=12, color=TEXT_DIM, bold=True)
add_text(slide, 6.5, y, 1.5, 0.3, "Count", font_size=12, color=TEXT_DIM, bold=True)
add_text(slide, 8, y, 1.5, 0.3, "% of All", font_size=12, color=TEXT_DIM, bold=True)

for i, (cwe, name, count, pct) in enumerate(cwe_data):
    y = 1.8 + i * 0.5
    add_text(slide, 1, y, 2, 0.4, cwe, font_size=15, color=RED, bold=True)
    add_text(slide, 3, y, 3, 0.4, name, font_size=15, color=TEXT_LIGHT)
    add_text(slide, 6.5, y, 1.5, 0.4, count, font_size=15, color=TEXT_WHITE, bold=True)
    add_text(slide, 8, y, 1.5, 0.4, pct, font_size=15, color=TEXT_DIM)

add_text(slide, 1, 5.8, 11, 0.5, "Key Insight: AI learned to avoid SQL injection (only 1 occurrence!) but systematically", font_size=16, color=AMBER)
add_text(slide, 1, 6.2, 11, 0.5, "fails at security configuration — headers, CSRF, rate limiting, cookie flags.", font_size=16, color=AMBER)

add_notes(slide, """JORDAN:
"Here's what our exploitation found. 18 CWE types, 2,295 total occurrences. The top 7 are shown here.

The dominant failure is CWE-693 — missing security headers. 62% of ALL findings. Every AI model systematically skips Content-Security-Policy, X-Frame-Options, and HSTS headers. It's the easiest thing to fix — one line of middleware — but no model does it.

CSRF is second at 271. No rate limiting is third at 144. XSS at 116.

The fascinating finding: SQL injection is almost ZERO. Only 1 occurrence in 4,505 tests. AI models have been trained on so much code with parameterized queries that they've actually learned to avoid string concatenation in SQL. That's a success story.

But for newer, less-documented vulnerabilities — CSRF, rate limiting, session management — they fail consistently."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 14: Post-Exploitation — Safety Prompts
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Jordan")
add_text(slide, 1, 0.5, 11, 0.6, "Post-Exploitation: Safety Prompt Impact", font_size=36, color=TEXT_WHITE, bold=True)
add_text(slide, 1, 1.2, 11, 0.5, "The single most important finding of our research", font_size=18, color=AMBER)

# Three columns
add_card(slide, 1, 2.0, 3.5, 2.2, "NONE PROMPT", "0.26%", "sec_pass@1 — 4 secure apps\n0.5% Sec(Working)\npass@1: 48.1%", RED)
add_card(slide, 5, 2.0, 3.5, 2.2, "GENERIC PROMPT", "0.07%", "sec_pass@1 — 1 secure app\n0.3% Sec(Working)\npass@1: 23.0%", RED)
add_card(slide, 9, 2.0, 3.5, 2.2, "SPECIFIC PROMPT", "9.73%", "sec_pass@1 — 141 secure apps\n70.5% Sec(Working)\npass@1: 13.8%", EMERALD)

add_text(slide, 1, 4.8, 11, 0.5, "Specific prompts: 37x improvement in sec_pass@1  |  141x improvement in Sec(Working)", font_size=18, color=EMERALD, bold=True)
add_text(slide, 1, 5.5, 11, 0.5, '"Write secure code" (generic) is WORSE than no prompt — 0.07% vs 0.26%', font_size=16, color=RED)
add_text(slide, 1, 6.0, 11, 0.5, "141 of 146 total secure apps come from the specific prompt", font_size=16, color=TEXT_LIGHT)

add_notes(slide, """JORDAN:
"This is our most important finding. We tested three safety prompt levels.

No prompt: just 'build this app.' Result: 0.26% sec_pass. Only 4 secure apps out of 1,533 tests. Only 1 model out of 15 produced ANY secure code.

Generic prompt: 'build this app, write secure code.' Result: 0.07% — actually WORSE. The vague instruction confuses the model. It tries to add security but does it wrong, breaking functionality without improving security.

Specific prompt: 'use bcrypt for passwords, set HttpOnly cookies, implement rate limiting.' Result: 9.73% — a 37x improvement. And among working code, 70.5% is secure — a 141x improvement.

141 of our 146 total secure apps come from the specific prompt. The AI knows HOW to write secure code. It just doesn't do it unless you tell it exactly what security measures to implement."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 15: Post-Exploitation — 3-Way Validation
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Jordan")
add_text(slide, 1, 0.5, 11, 0.6, "Post-Exploitation: 3-Way Validation", font_size=36, color=TEXT_WHITE, bold=True)

# Comparison table header
y = 1.5
add_text(slide, 1, y, 3, 0.3, "Vulnerability", font_size=14, color=TEXT_DIM, bold=True)
add_text(slide, 4.5, y, 2, 0.3, "CodeStrike", font_size=14, color=PURPLE, bold=True)
add_text(slide, 7, y, 2, 0.3, "ZAP", font_size=14, color=AMBER, bold=True)
add_text(slide, 9.5, y, 2, 0.3, "Manual", font_size=14, color=EMERALD, bold=True)

vulns = [
    ("Missing Headers", "80%", "Found", "Found"),
    ("CSRF", "100%", "Can't", "Found"),
    ("Brute Force", "25%", "Can't", "Found"),
    ("Hardcoded Secrets", "0%", "Can't", "Found"),
    ("Access Control", "0%", "Can't", "Found"),
    ("SSRF", "0%", "Can't", "Found"),
    ("Business Logic", "0%", "Can't", "Found"),
]

for i, (vuln, cs, zap, manual) in enumerate(vulns):
    y = 2.0 + i * 0.45
    add_text(slide, 1, y, 3, 0.4, vuln, font_size=14, color=TEXT_LIGHT)
    cs_color = EMERALD if cs in ["80%", "100%"] else AMBER if cs == "25%" else RED
    add_text(slide, 4.5, y, 2, 0.4, cs, font_size=14, color=cs_color, bold=True)
    zap_color = EMERALD if zap == "Found" else RED
    add_text(slide, 7, y, 2, 0.4, zap, font_size=14, color=zap_color)
    add_text(slide, 9.5, y, 2, 0.4, manual, font_size=14, color=EMERALD)

add_text(slide, 1, 5.5, 5, 0.4, "CodeStrike: 100% precision, 27% recall", font_size=16, color=PURPLE, bold=True)
add_text(slide, 7, 5.5, 5, 0.4, "ZAP: 14.3% agreement", font_size=16, color=AMBER, bold=True)
add_text(slide, 1, 6.2, 11, 0.5, "Why CodeStrike beats ZAP: it runs INSIDE Docker — can access DB, test rate limiting, understand auth logic", font_size=15, color=TEXT_LIGHT)

add_notes(slide, """JORDAN:
"Here's our 3-way comparison. This table shows detection rates per vulnerability type.

CodeStrike catches headers at 80%, CSRF at 100% when present, and brute force at 25%. But it can't find hardcoded secrets, access control issues, or business logic — those require human reasoning.

ZAP can only detect headers. It can't test brute force, can't access the database, can't understand authentication logic. Its XSS scanner injects script tags but our APIs return JSON, so it never detects reflection.

Manual testing found everything — but it doesn't scale. That's the trade-off.

CodeStrike's architectural advantage: it runs INSIDE the Docker container. It can query SQLite to check if passwords are hashed, send 15 rapid requests to test rate limiting, and create two users to test IDOR. ZAP can only send HTTP requests from outside."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 16: The Bug We Found
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "Jordan")
add_text(slide, 1, 0.5, 11, 0.6, "Critical Discovery: The 'Secure by Crash' Bug", font_size=36, color=TEXT_WHITE, bold=True)

add_text(slide, 1, 1.5, 11, 0.5, "The original benchmark had a measurement bug that inflated security rates by 20x", font_size=18, color=RED)

add_card(slide, 1, 2.5, 5, 1.5, "BEFORE (BUG)", "~94%", "sonnet-4.5-thinking sec_pass\nCrashed code counted as 'secure'\nbecause no CWEs were detected on broken code", RED)
add_card(slide, 7, 2.5, 5, 1.5, "AFTER (FIXED)", "4.4%", "sonnet-4.5-thinking sec_pass\nNow requires functional_pass = true\nCrashed code is NOT secure", EMERALD)

add_text(slide, 1, 4.5, 11, 0.4, "Our Fix: Three Metrics", font_size=20, color=AMBER, bold=True)
add_bullet_list(slide, 1, 5.0, 11, 2, [
    ("sec_pass@1 (3.1%) = functional + zero CWEs / total tests", AMBER),
    ("true_sec@1 (0.7%) = functional + zero CWEs + zero test crashes / total tests", EMERALD),
    ("Sec(Working) (11.3%) = zero CWEs / functional tests only — fairest comparison", PURPLE),
], font_size=16)

add_notes(slide, """JORDAN:
"We discovered a critical measurement bug. The original benchmark counted crashed apps as 'secure' because if the code crashes, no security tests run, so no CWEs are detected. Zero CWEs looked like 'secure.'

This made sonnet-4.5-thinking appear to have a 94% security rate. The truth? 93% of its code crashed. When we fixed the formula to require functional_pass = true, the real rate dropped to 4.4%.

We introduced three metrics to address this: sec_pass@1 requires the code to work AND have zero CWEs. true_sec@1 additionally requires all security tests to run cleanly. And Sec(Working) divides only by functional apps — the fairest comparison between models.

This is a methodological contribution. Any future benchmark should use these metrics."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 17: Summary — Key Findings
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "[Member 4]")
add_text(slide, 1, 0.5, 11, 0.6, "Summary: Key Findings", font_size=36, color=TEXT_WHITE, bold=True)

findings = [
    ("Less than 1% of AI code is truly secure", "true_sec@1 = 0.7% — 31 of 4,505 apps", RED),
    ("71.4% of AI code doesn't even work", "Code quality is the primary problem", RED),
    ("Safety prompts give 37-141x improvement", "0.26% -> 9.73% sec_pass, 0.5% -> 70.5% Sec(Working)", EMERALD),
    ("'Write secure code' is useless", "Generic prompt: 0.07% — worse than no prompt", AMBER),
    ("Thinking mode doesn't help security", "+0.3 pp average — negligible improvement", AMBER),
    ("Express is the only safe framework", "144/146 secure apps — Flask 2, Go-Fiber 0", BLUE),
    ("AI avoids SQLi but fails at configuration", "1 SQLi vs 1,427 missing headers", PURPLE),
    ("7/10 OWASP categories have findings", "Security Misconfiguration dominates (A05)", RED),
]

for i, (finding, detail, color) in enumerate(findings):
    y = 1.3 + i * 0.7
    add_text(slide, 1, y, 7, 0.35, finding, font_size=16, color=color, bold=True)
    add_text(slide, 1, y + 0.3, 7, 0.3, detail, font_size=12, color=TEXT_DIM)

add_notes(slide, """[MEMBER 4] (Speaker 4):
"Thanks Jordan. Let me summarize our key findings.

Number 1: Less than 1% of AI-generated code is truly secure. That's 31 apps out of 4,505.

Number 2: Most AI code doesn't even work. 71% crashes before we can test security.

Number 3: This is the actionable finding — specific safety prompts improve security by 37x to 141x. Just telling the AI exactly what security measures to implement makes a massive difference.

Number 4: Vague instructions like 'write secure code' are useless — actually worse than nothing.

Number 5: Thinking mode, where the AI reasons before coding, doesn't help security at all.

Number 6: Express is the only framework where AI consistently writes secure code. 144 of 146 secure apps.

Number 7: AI learned to avoid SQL injection but fails at basic configuration.

Number 8: We found vulnerabilities in 7 of 10 OWASP categories."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 18: Recommendations
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_speaker_tag(slide, "[Member 4]")
add_text(slide, 1, 0.5, 11, 0.6, "Recommendations & Remediation", font_size=36, color=TEXT_WHITE, bold=True)

add_text(slide, 1, 1.3, 5.5, 0.4, "For Developers Using AI Code", font_size=20, color=EMERALD, bold=True)
add_bullet_list(slide, 1, 1.8, 5.5, 2.8, [
    "Always use specific safety prompts — list every security requirement explicitly",
    "Never trust AI code without security review",
    "Use Express/Node.js for best security outcomes",
    "Add security middleware (helmet.js, flask-talisman) to every project",
    "Implement rate limiting separately — AI never adds it",
    "Run automated security scanners (CodeStrike, ZAP) in CI/CD",
], font_size=15)

add_text(slide, 7, 1.3, 5.5, 0.4, "For AI Model Providers", font_size=20, color=PURPLE, bold=True)
add_bullet_list(slide, 7, 1.8, 5.5, 2.8, [
    "Train models on security middleware patterns, not just secure coding",
    "Default to adding security headers in all generated web code",
    "Include CSRF protection by default in form-handling code",
    "Add rate limiting patterns to authentication code generation",
    "Improve Go/Python security training data (Express is 100x better)",
], font_size=15)

add_text(slide, 1, 5.0, 11, 0.4, "For Security Teams", font_size=20, color=AMBER, bold=True)
add_bullet_list(slide, 1, 5.5, 11, 1.5, [
    "Treat ALL AI-generated code as untrusted — same as third-party code",
    "Layer your testing: automated (CodeStrike) for scale + manual for depth + DAST (ZAP) for config",
    "Use true_sec@1, not sec_pass@1, to avoid the 'secure by crash' trap",
], font_size=15)

add_notes(slide, """[MEMBER 4]:
"Based on our findings, here are our recommendations.

For developers: Always use specific safety prompts. Don't say 'write secure code' — list every requirement: 'use bcrypt, set HttpOnly, implement rate limiting.' This single change gives you a 141x security improvement.

For AI providers: The training gap is clear. Models know parameterized SQL queries but not security headers. Train on middleware patterns, not just coding patterns.

For security teams: Treat AI code like untrusted third-party code. Use layered testing. And use true_sec@1 as your metric — not the inflated sec_pass@1.

Now let me show you all of this live in our dashboard."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 19: Dashboard Demo Transition
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1, 2, 11, 1.2, "Live Dashboard Demo", font_size=48, color=EMERALD, bold=True)
add_text(slide, 1, 3.5, 11, 0.6, "dashboard-wheat-iota-87.vercel.app", font_size=24, color=BLUE)
add_text(slide, 1, 4.5, 11, 0.5, "5 Pages: Overview  |  Models  |  Vulnerabilities  |  Compare  |  Pentest", font_size=18, color=TEXT_LIGHT)

add_notes(slide, """[MEMBER 4]:
"Now let me switch to our live dashboard. I'll walk you through each page.

Demo flow:
1. OVERVIEW: Show the Security Funnel (click through 6 steps). Point out stat cards. Show the Model Ranking table — toggle safety prompts from None (wall of zeros) to Specific (numbers jump).

2. MODELS: Click on sonnet-4-standard. Show the radar chart, top vulnerabilities, and the All Results table. Point out the 4 security states: crashed, vulnerable, uncertain, secure.

3. VULNERABILITIES: Show the OWASP Top 10 2025 section — expand A03 Injection and A05 Misconfiguration. Show the CWE treemap.

4. COMPARE: Go to Safety Prompts tab. Toggle the chart. Show the table with all 3 columns (None/Generic/Specific). Switch to Frameworks tab — toggle from sec_pass to pass@1 to show all bars.

5. PENTEST: Show the 3-way comparison table. Show the 'Why CodeStrike Catches More' section."
""")

# ═══════════════════════════════════════════════════════════
# SLIDE 20: Thank You
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1, 2, 11, 1, "Thank You", font_size=52, color=TEXT_WHITE, bold=True)
add_text(slide, 1, 3.2, 11, 0.5, "Questions?", font_size=28, color=EMERALD)
add_text(slide, 1, 4.5, 11, 0.4, "CodeStrike Security Dashboard: dashboard-wheat-iota-87.vercel.app", font_size=16, color=BLUE)
add_text(slide, 1, 5.0, 11, 0.4, "GitHub: github.com/iyassh/baxbench-extended", font_size=16, color=TEXT_DIM)
add_text(slide, 1, 5.5, 11, 0.4, "COMP 4210 Ethical Hacking  |  Group 8  |  April 2026", font_size=14, color=TEXT_DIM)

add_notes(slide, """ALL:
"Thank you for listening. We're happy to take questions.

Expected questions and answers:

Q: Why is the security rate so low?
A: Two reasons — 71% crashes + 89% of working code has vulnerabilities. AI focuses on functionality, not security.

Q: Isn't 4.4% too low to be useful?
A: With specific safety prompts, 70.5% of working code is secure. The key is HOW you prompt.

Q: How do you know CodeStrike is accurate?
A: 100% precision — zero false positives across 10 manually pentested apps.

Q: What about the 73% recall gap?
A: Expected for automated tools. Hardcoded creds, IDOR, business logic require human reasoning.

Q: Is thinking mode worth the extra cost?
A: No, for security. +0.3 pp is negligible.

Q: Main recommendation?
A: Use specific safety prompts. List every security requirement explicitly. 141x improvement."
""")

# Save
output_path = "/Users/yassh/baxbench/docs/CodeStrike_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
